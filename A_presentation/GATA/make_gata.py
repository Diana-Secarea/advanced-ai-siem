#!/usr/bin/env python3
"""Build the GATA short defense deck (11 slides + labeled appendix).

Follows the academic-pptx skill: white content slides, action titles,
one exhibit per slide, in-slide citations, conclusions last (stays up in Q&A).

Run:  ../../ai_threat_engine_starter/venv/bin/python make_gata.py
Outputs: GATA_Defense_Presentation.pptx  (+ light-theme charts in img/)
"""
import os
import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
import numpy as np
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")
os.makedirs(IMG, exist_ok=True)

# ---------- academic palette (skill defaults + one warm emphasis) ----------
NAVY   = RGBColor(0x1F, 0x4E, 0x79)   # primary — titles, dark slides
BLUE   = RGBColor(0x2E, 0x75, 0xB6)   # accent — headers, focal series
SKY    = RGBColor(0xEB, 0xF3, 0xFA)   # callout fill
BODY   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x77, 0x77, 0x77)
RULE   = RGBColor(0xCC, 0xCC, 0xCC)
HL_BG  = RGBColor(0xFF, 0xF2, 0xCC)   # highlight box fill
HL_BR  = RGBColor(0xE6, 0xC8, 0x00)
HL_TX  = RGBColor(0x7A, 0x52, 0x00)
GREEN  = RGBColor(0x2E, 0x7D, 0x5B)   # semantic good
REDX   = RGBColor(0xB0, 0x35, 0x2F)   # semantic alert (sparing)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTB = RGBColor(0xA0, 0xBB, 0xDD)   # light blue text on navy
FONT = "Calibri"

# hex strings for matplotlib
H_NAVY, H_BLUE, H_MUT, H_BODY = "#1F4E79", "#2E75B6", "#777777", "#2D2D2D"
H_AMBER, H_RED, H_GREEN, H_GRID = "#E6A817", "#B0352F", "#2E7D5B", "#E2E8F0"

plt.rcParams.update({
    "font.family": "DejaVu Sans", "figure.facecolor": "white",
    "axes.facecolor": "white", "axes.edgecolor": H_GRID,
    "axes.labelcolor": H_MUT, "xtick.color": H_MUT, "ytick.color": H_BODY,
    "font.size": 15,
})

# ============================ CHARTS (light theme) ============================

def chart_ablation():
    feats = ["rule_id  (f10)", "unknown_user_flag  (f15)", "rule_level  (f9)",
             "failed_count  (f3)", "privileged_account_change  (f16)"]
    df1 = [-2.8, -3.5, -4.1, -5.7, -6.3]
    fig, ax = plt.subplots(figsize=(8.4, 4.4), dpi=200)
    colors = [H_BLUE] * 4 + [H_NAVY]
    bars = ax.barh(feats, df1, color=colors, height=0.62)
    for b, v in zip(bars, df1):
        ax.text(v - 0.12, b.get_y() + b.get_height() / 2, f"{v:.1f}%",
                va="center", ha="right", fontsize=15, color=H_BODY, fontweight="bold")
    ax.set_xlabel("Δ F1-score when the feature is removed (pp)", fontsize=14)
    ax.set_xlim(-7.6, 0)
    ax.axvline(0, color=H_GRID, lw=1)
    ax.spines[["top", "right"]].set_visible(False)
    ax.xaxis.grid(True, color=H_GRID, lw=0.6)
    ax.set_axisbelow(True)
    # annotate the key finding
    ax.annotate("strongest single signal", xy=(-6.3, 4), xytext=(-4.6, 4.32),
                fontsize=14, color=H_NAVY, fontweight="bold",
                arrowprops=dict(arrowstyle="->", color=H_NAVY, lw=1.4))
    plt.tight_layout()
    fig.savefig(os.path.join(IMG, "ablation_light.png"))
    plt.close(fig)


def chart_confusion():
    fig, ax = plt.subplots(figsize=(6.6, 5.0), dpi=200)
    cells = [("TP", 596, H_NAVY, "white"), ("FN", 41, "#F5DFA6", H_BODY),
             ("FP", 130, "#F5DFA6", H_BODY), ("TN", 1404, H_BLUE, "white")]
    pos = {(0, 1): 0, (1, 1): 1, (0, 0): 2, (1, 0): 3}
    for (cx, cy), idx in pos.items():
        lab, n, fc, tc = cells[idx]
        ax.add_patch(plt.Rectangle((cx, cy), 0.96, 0.96, color=fc))
        ax.text(cx + 0.48, cy + 0.58, f"{n:,}", ha="center", va="center",
                fontsize=30, fontweight="bold", color=tc)
        ax.text(cx + 0.48, cy + 0.26, lab, ha="center", va="center",
                fontsize=16, color=tc)
    ax.set_xlim(-0.35, 2.05); ax.set_ylim(-0.32, 2.28)
    ax.text(0.48, 2.06, "flagged anomaly", ha="center", fontsize=14, color=H_MUT)
    ax.text(1.48, 2.06, "flagged normal", ha="center", fontsize=14, color=H_MUT)
    ax.text(-0.18, 1.48, "attack", va="center", rotation=90, fontsize=14, color=H_MUT)
    ax.text(-0.18, 0.48, "clean", va="center", rotation=90, fontsize=14, color=H_MUT)
    # focal annotation on FN
    ax.annotate("only 41 missed\n(IF alone: 62 · AE alone: 208)",
                xy=(1.5, 1.5), xytext=(1.06, -0.3),
                fontsize=14, color=H_RED, fontweight="bold",
                arrowprops=dict(arrowstyle="->", color=H_RED, lw=1.5))
    ax.axis("off")
    plt.tight_layout()
    fig.savefig(os.path.join(IMG, "confusion_light.png"))
    plt.close(fig)


def chart_rag_generation():
    metrics = ["Faithfulness\n(LLM-judge)", "Refusal accuracy\n(unanswerable set)",
               "Citation coverage", "Citation precision"]
    vals = [71.0, 83.3, 88.6, 98.3]
    fig, ax = plt.subplots(figsize=(8.4, 4.3), dpi=200)
    colors = [H_MUT, H_BLUE, H_BLUE, H_NAVY]
    bars = ax.barh(metrics, vals, color=colors, height=0.6)
    for b, v in zip(bars, vals):
        ax.text(v - 1.2, b.get_y() + b.get_height() / 2, f"{v:.1f}%",
                va="center", ha="right", fontsize=15, fontweight="bold",
                color="white" if v > 75 else "white")
    ax.set_xlim(0, 100)
    ax.set_xlabel("Score (%) — 82 evaluation queries + 12 unanswerable probes", fontsize=13)
    ax.spines[["top", "right"]].set_visible(False)
    ax.xaxis.grid(True, color=H_GRID, lw=0.6)
    ax.set_axisbelow(True)
    ax.annotate("117 of 119 citations verified", xy=(98.3, 3), xytext=(52, 3.38),
                fontsize=14, color=H_NAVY, fontweight="bold",
                arrowprops=dict(arrowstyle="->", color=H_NAVY, lw=1.4))
    plt.tight_layout()
    fig.savefig(os.path.join(IMG, "rag_gen_light.png"))
    plt.close(fig)


def chart_rag_retrieval():
    cats = ["MITRE ATT&CK\ntechniques", "Wazuh rules\n& responses", "C2 / exfiltration",
            "Lateral movement", "CVE / vuln.\ndetails", "APT group\nprofiles"]
    recall = [46.3, 71.4, 80.2, 85.7, 100.0, 100.0]
    qn_ = [22, 7, 14, 12, 18, 9]
    # focal: MITRE weakest (amber-red), CVE & APT strongest (navy), rest blue
    colors = [H_RED, H_BLUE, H_BLUE, H_BLUE, H_NAVY, H_NAVY]
    fig, ax = plt.subplots(figsize=(8.3, 4.6), dpi=200)
    bars = ax.barh(cats, recall, color=colors, height=0.66)
    for b, v, q in zip(bars, recall, qn_):
        ax.text(v - 1.5, b.get_y() + b.get_height() / 2, f"{v:.0f}%",
                va="center", ha="right", fontsize=15, fontweight="bold",
                color="white")
        ax.text(101, b.get_y() + b.get_height() / 2, f"n={q}",
                va="center", ha="left", fontsize=12, color=H_MUT)
    ax.set_xlim(0, 100)
    ax.set_xlabel("Recall@5 (%) — fraction of relevant chunks retrieved in the top 5", fontsize=13)
    ax.spines[["top", "right"]].set_visible(False)
    ax.xaxis.grid(True, color=H_GRID, lw=0.6)
    ax.set_axisbelow(True)
    ax.annotate("weak spot", xy=(46.3, 0), xytext=(66, 0.02),
                fontsize=13, color=H_RED, fontweight="bold", va="center",
                arrowprops=dict(arrowstyle="->", color=H_RED, lw=1.4))
    plt.tight_layout()
    fig.savefig(os.path.join(IMG, "rag_retrieval_light.png"))
    plt.close(fig)


# ============================ PPTX helpers ============================

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
M = Inches(0.6)          # margin
CW = Inches(12.13)       # content width

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]


def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY if dark else WHITE
    return s


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def tx(slide, x, y, w, h, text, size=20, color=BODY, bold=False, italic=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT
    return tb


def action_title(slide, text, size=25):
    tx(slide, M, Inches(0.28), CW, Inches(1.0), text, size=size, color=NAVY, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(1.26), CW, Emu(22860))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background()


def rich_bullets(slide, x, y, w, h, items, size=19, gap=10, marker_color=BLUE):
    """items: list of (bold_lead, rest) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        rm = p.add_run(); rm.text = "▪  "
        rm.font.size = Pt(size); rm.font.color.rgb = marker_color; rm.font.name = FONT
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = BODY; r1.font.name = FONT
        if rest:
            r2 = p.add_run(); r2.text = " " + rest
            r2.font.size = Pt(size); r2.font.color.rgb = BODY; r2.font.name = FONT
    return tb


def cite(slide, text, y=Inches(7.02)):
    tx(slide, M, y, CW, Inches(0.36), text, size=12, color=MUTED)


def col_header(slide, x, y, w, text, color=BLUE):
    tx(slide, x, y, w, Inches(0.4), text, size=20, color=color, bold=True)


def box(slide, x, y, w, h, fill, line_color=None, line_w=1.0, shadow=False,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.12
    except Exception:
        pass
    return sh


def box_text(sh, lines, size=14, color=BODY, bold_first=True):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(size if i == 0 else size - 2)
        r.font.bold = bold_first and i == 0
        r.font.color.rgb = color; r.font.name = FONT


def arrow(slide, x1, y1, x2, y2, color=MUTED, w=2.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)  # straight
    conn.line.color.rgb = color; conn.line.width = Pt(w)
    le = conn.line._get_or_add_ln()
    tail = le.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    le.append(tail)


def highlight_box(slide, x, y, w, h, text, size=14):
    sh = box(slide, x, y, w, h, HL_BG, HL_BR, 1.0)
    box_text(sh, [text], size=size, color=HL_TX, bold_first=True)
    return sh


def pic(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(os.path.join(IMG, path), x, y, w, h)


# ============================ SLIDES ============================

def s01_title():
    s = new_slide(dark=True)
    tx(s, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.5),
       "BACHELOR THESIS DEFENSE", size=15, color=LIGHTB)
    tx(s, Inches(0.9), Inches(1.5), Inches(11.6), Inches(2.6),
       "Architecture and Evaluation of a SIEM\nIntegrating Unsupervised Anomaly Detection\nand Knowledge-Grounded Language Models",
       size=33, color=WHITE, bold=True, spacing=1.08)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.35), Inches(2.2), Emu(45720))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    tx(s, Inches(0.9), Inches(4.6), Inches(11), Inches(0.45),
       "Secărea Diana Maria", size=21, color=WHITE, bold=True)
    tx(s, Inches(0.9), Inches(5.12), Inches(11.6), Inches(1.0),
       "Coordinator: prof. univ. dr. Cătălin Boja\nBucharest University of Economic Studies · Faculty of Cybernetics, Statistics and Economic Informatics",
       size=15, color=LIGHTB)
    tx(s, Inches(0.9), Inches(6.45), Inches(11), Inches(0.4),
       "Economic Informatics Program · Bucharest 2026", size=13, color=LIGHTB)
    add_notes(s, "Good morning. I will present my thesis: an AI threat-hunting engine built on the "
                 "Wazuh SIEM — unsupervised anomaly detection plus a knowledge-grounded LLM copilot. "
                 "Everything I show is implemented, deployed and evaluated; there will also be a short live demo.")


def s02_motivation():
    s = new_slide()
    action_title(s, "Rule-based SIEMs leave analysts drowning in alerts while novel attacks pass undetected")
    rich_bullets(s, M, Inches(1.55), Inches(7.1), Inches(4.4), [
        ("Alert fatigue is structural:", "SOCs average ~4,490 alerts/day; ~67% are never investigated for lack of analyst bandwidth."),
        ("Rules only match the known.", "A zero-day or custom tool produces events no rule ID covers — detection is reactive by design."),
        ("Static thresholds misfire.", "“8 failed logins in 120 s” fits no specific environment → false-positive floods."),
        ("Alerts arrive context-free.", "“Possible kernel rootkit (rule 521)” still requires manual research into T1014, indicators, remediation."),
    ], size=19, gap=14)
    # right: the question box
    qb = box(s, Inches(8.1), Inches(1.7), Inches(4.55), Inches(3.5), SKY, BLUE, 1.5)
    box_text(qb, ["Research question",
                  "",
                  "Can unsupervised anomaly scoring",
                  "+ retrieval-grounded reasoning",
                  "reduce triage time and improve",
                  "analyst trust — on a real,",
                  "running Wazuh deployment?"], size=18, color=NAVY)
    highlight_box(s, Inches(8.1), Inches(5.45), Inches(4.55), Inches(0.62),
                  "Approach: score every alert, explain every score", size=14)
    cite(s, "Industry SOC statistics [7]; WEF Global Risks Report 2024 [4]; Wazuh documentation [2] — thesis §2")
    add_notes(s, "Three problems motivate the work: alert fatigue (most alerts never get looked at), "
                 "the structural blindness of signature detection to novel attacks, and the missing context "
                 "when an alert does fire. My research question: can unsupervised ML plus grounded RAG fix "
                 "triage on a real deployment — not a toy dataset.")


def s03_system():
    s = new_slide()
    action_title(s, "The AI Threat-Hunting Engine extends Wazuh end-to-end: score → explain → keep knowledge fresh")
    p = pic(s, "architecture_thesis.png", M, Inches(1.5), w=Inches(7.55))
    # right rail: the three contributions
    x = Inches(8.45); w = Inches(4.25)
    labels = [
        ("1 · Anomaly ensemble", "Isolation Forest + Autoencoder score every alert 0–100, fused into confidence tiers", NAVY),
        ("2 · RAG copilot", "local LLM grounded in MITRE ATT&CK, YARA, Sigma, CVEs — cited answers in the dashboard", BLUE),
        ("3 · CVE ingestion agent", "scheduled agent fetches, scores and routes new CVEs; PostgreSQL keeps the audit ledger", GREEN),
    ]
    y = Inches(1.75)
    for head, body, c in labels:
        b = box(s, x, y, w, Inches(1.42), SKY, c, 1.4)
        tf = b.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        r = p1.add_run(); r.text = head
        r.font.size = Pt(17); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(13.5); r2.font.color.rgb = BODY; r2.font.name = FONT
        y += Inches(1.62)
    cite(s, "Figure 3.1.1 — Architecture of the proposed system (thesis, p. 9). Endpoint & manager layers = Wazuh; engine, storage, visualisation layers = this work.")
    add_notes(s, "This is the architecture figure from the thesis. Left half is standard Wazuh: agents, "
                 "manager, rule engine. Everything else is my contribution — the AI Threat Hunting Engine "
                 "with the ensemble scorer, the RAG system with LLM, the CVE collector agent, a storage "
                 "layer with Qdrant vectors + PostgreSQL audit, and the visualisation layer with dashboards. "
                 "Three cooperating layers: score, explain, stay current.")


def s04_technologies():
    s = new_slide()
    action_title(s, "A fully local, open-source stack: security data never leaves the machine")
    cols = [
        ("Detection (ML)", NAVY, [
            "scikit-learn — Isolation Forest",
            "MLP Autoencoder 16→8→4→8→16",
            "NumPy · pandas · joblib artefacts",
            "16-feature extractor (custom)",
        ]),
        ("Explanation (RAG)", BLUE, [
            "Qdrant vector DB — HNSW",
            "all-MiniLM-L6-v2 embeddings (384-d)",
            "BM25 sparse (fastembed) + RRF",
            "Ollama · llama3.2 — local LLM",
        ]),
        ("Platform", GREEN, [
            "Wazuh 4.14.1 SIEM (fork base)",
            "Flask API + dashboards :5000",
            "PostgreSQL 17 — audit ledger",
            "Docker Compose · Python 3.13 · WSL2",
        ]),
    ]
    x = M
    cw = Inches(3.9)
    for head, c, items in cols:
        hdr = box(s, x, Inches(1.6), cw, Inches(0.52), c)
        box_text(hdr, [head], size=17, color=WHITE)
        bx = box(s, x, Inches(2.12), cw, Inches(3.3), SKY, c, 1.0,
                 shape=MSO_SHAPE.RECTANGLE)
        tf = bx.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.14); tf.margin_top = Inches(0.12)
        for i, it in enumerate(items):
            pgh = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            pgh.space_after = Pt(9)
            r = pgh.add_run(); r.text = "▪ " + it
            r.font.size = Pt(15.5); r.font.color.rgb = BODY; r.font.name = FONT
        x += cw + Inches(0.22)
    highlight_box(s, Inches(2.35), Inches(5.75), Inches(8.6), Inches(0.68),
                  "Why local? Alerts contain confidential system data — no third-party API, no per-token cost", size=15)
    cite(s, "Thesis §4.1 Technologies Used")
    add_notes(s, "Deliberately boring, deliberately local stack. Detection: scikit-learn ensemble. "
                 "Explanation: Qdrant with hybrid dense+sparse retrieval and a local llama3.2 through Ollama — "
                 "because SIEM alerts are confidential, nothing goes to external APIs. Platform: fork of Wazuh, "
                 "Flask dashboards, Postgres for auditability, all in Docker on one Linux box.")


def s05_ensemble():
    s = new_slide()
    action_title(s, "Two models that fail differently: extremes get isolated, unusual combinations fail to reconstruct")
    # pipeline diagram (native shapes)
    y0 = Inches(1.75)
    fx = M
    fb = box(s, fx, y0 + Inches(0.85), Inches(2.15), Inches(1.15), SKY, NAVY, 1.2)
    box_text(fb, ["Wazuh alert", "16 features", "z-scaled"], size=15, color=NAVY)
    ifb = box(s, Inches(3.6), y0, Inches(2.9), Inches(1.25), SKY, NAVY, 1.4)
    box_text(ifb, ["Isolation Forest", "random cuts isolate", "extreme features"], size=15, color=NAVY)
    aeb = box(s, Inches(3.6), y0 + Inches(1.75), Inches(2.9), Inches(1.25), SKY, BLUE, 1.4)
    box_text(aeb, ["Autoencoder 16→8→4→8→16", "reconstruction error flags", "odd combinations"], size=15, color=BLUE)
    fu = box(s, Inches(7.3), y0 + Inches(0.85), Inches(2.3), Inches(1.15), NAVY)
    box_text(fu, ["Fusion", "0.45·IF + 0.55·AE", "calibrated 0–100"], size=15, color=WHITE)
    arrow(s, fx + Inches(2.15), y0 + Inches(1.42), Inches(3.6), y0 + Inches(0.62))
    arrow(s, fx + Inches(2.15), y0 + Inches(1.42), Inches(3.6), y0 + Inches(2.37))
    arrow(s, Inches(6.5), y0 + Inches(0.62), Inches(7.3), y0 + Inches(1.3))
    arrow(s, Inches(6.5), y0 + Inches(2.37), Inches(7.3), y0 + Inches(1.6))
    # tier chips
    tiers = [("CRITICAL", "both flag", REDX), ("HIGH", "only AE", RGBColor(0xC5, 0x7A, 0x00)),
             ("POSSIBLE", "only IF ×0.65", BLUE), ("NORMAL", "neither", GREEN)]
    ty = y0 + Inches(0.1)
    for name, sub, c in tiers:
        t = box(s, Inches(10.1), ty, Inches(2.55), Inches(0.66), WHITE, c, 1.6)
        tf = t.text_frame; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        pgh = tf.paragraphs[0]; pgh.alignment = PP_ALIGN.CENTER
        r = pgh.add_run(); r.text = name + "  ·  " + sub
        r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
        ty += Inches(0.78)
    arrow(s, Inches(9.6), y0 + Inches(1.42), Inches(10.1), y0 + Inches(1.42))
    # adversarial argument
    rich_bullets(s, M, Inches(5.05), Inches(12.1), Inches(1.6), [
        ("Evading both is contradictory:", "blending into density (fools IF) makes you statistically average — easy to reconstruct; mimicking the normal manifold (fools AE) needs common feature values — easy to isolate."),
        ("Tiers carry meaning, not just a number:", "agreement = act now; single-model flags = review when time allows. 91.3% of detected attacks land in CRITICAL."),
    ], size=17, gap=8)
    cite(s, "Isolation Forest: Liu, Ting & Zhou (2008) [20]; autoencoders: Michelucci (2022) [23] — thesis §3.2")
    add_notes(s, "The core design: one alert, sixteen features, two unsupervised models with opposite failure "
                 "modes. IF isolates extreme single features fast; the AE learns to reconstruct normal alerts, "
                 "so unusual combinations reconstruct badly even when no single feature is extreme. Fusion is "
                 "weighted 45/55 with tier logic from the binary decisions. Key defense point: the evasion "
                 "strategies against the two models contradict each other — you cannot be statistically average "
                 "and statistically common at the same time.")


def s06_features():
    s = new_slide()
    action_title(s, "All 16 features are ablation-tested: privilege changes, failed logins and rule severity carry the signal")
    pic(s, "ablation_light.png", M, Inches(1.55), w=Inches(7.3))
    col_header(s, Inches(8.3), Inches(1.55), Inches(4.4), "What the ablation shows")
    rich_bullets(s, Inches(8.3), Inches(2.05), Inches(4.4), Inches(3.2), [
        ("Retrain from scratch per variant", "(< 5 s on ~2,200 alerts) — not approximated."),
        ("privileged_account_change:", "without it, new-user/promiscuous-mode attacks drop to 0% detection."),
        ("failed_count:", "the brute-force burst signal; CIS-aware zeroing avoids compliance-scan false positives."),
    ], size=16, gap=10)
    highlight_box(s, Inches(8.3), Inches(5.35), Inches(4.4), Inches(1.15),
                  "Counter-example: adding mitre_count → +8.9% false positives (training/production drift) — removed", size=14)
    cite(s, "Ablation study — thesis §3.2 and Appendix 5 (baseline: 16 features, F1 = 89.2%, FPR = 8.0%)")
    add_notes(s, "Feature engineering was validated, not guessed: remove one feature, retrain fully, measure. "
                 "Privileged-account-change is the strongest signal — attacker persistence. The honest part "
                 "examiners like: a feature that looked obviously useful, MITRE technique count, made the model "
                 "worse — Wazuh tags benign events with MITRE in production but not in my training data, so I "
                 "removed it. Distribution shift, addressed empirically.")


def s07_training():
    s = new_slide()
    action_title(s, "Zero labels, contaminated data: iterative trimming makes ~30% attack contamination trainable")
    # left: trimming pipeline
    y = Inches(1.8)
    steps = [
        ("Round 1", "train AE on ALL 2,171 alerts → drop top 25% reconstruction errors", NAVY),
        ("Round 2", "retrain on the remainder → drop top 15% errors", BLUE),
        ("Round 3", "final fit on the clean estimate → calibrate on 2nd–98th percentiles → 0–100", GREEN),
    ]
    for name, body, c in steps:
        b = box(s, M, y, Inches(6.9), Inches(1.0), SKY, c, 1.3)
        tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        tf.margin_left = Inches(0.14)
        pgh = tf.paragraphs[0]
        r = pgh.add_run(); r.text = name + " — "
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
        r2 = pgh.add_run(); r2.text = body
        r2.font.size = Pt(15.5); r2.font.color.rgb = BODY; r2.font.name = FONT
        if y < Inches(4):
            arrow(s, Inches(4.0), y + Inches(1.0), Inches(4.0), y + Inches(1.28), MUTED, 1.8)
        y += Inches(1.28)
    tx(s, M, Inches(5.75), Inches(6.9), Inches(0.9),
       "Rationale: attacks reconstruct badly, so the worst-reconstructed tail is where they hide. The model bootstraps its own clean set — no analyst labelling, self-correcting on retrain.",
       size=14.5, color=MUTED, italic=True)
    # right: the honest trade table
    col_header(s, Inches(7.95), Inches(1.55), Inches(4.7), "Baseline vs. label-free")
    rows = [("", "clean-only (baseline)", "unsupervised (final)"),
            ("False-positive rate", "8.83%", "6.2%  ▼"),
            ("Precision", "83.7%", "85.8%  ▲"),
            ("Recall", "95.7%", "90.3%"),
            ("F1-score", "89.3%", "88.0%")]
    ty = Inches(2.15)
    for i, (a, b_, c_) in enumerate(rows):
        for j, (val, xx, ww) in enumerate([(a, Inches(7.95), Inches(1.75)),
                                           (b_, Inches(9.7), Inches(1.5)),
                                           (c_, Inches(11.2), Inches(1.55))]):
            cell = box(s, xx, ty, ww, Inches(0.5),
                       NAVY if i == 0 else (WHITE if j == 0 else SKY),
                       RULE, 0.75, shape=MSO_SHAPE.RECTANGLE)
            box_text(cell, [val], size=13 if i == 0 else 14,
                     color=WHITE if i == 0 else (BODY if j == 0 else NAVY),
                     bold_first=(i == 0 or j == 0))
        ty += Inches(0.5)
    tx(s, Inches(7.95), Inches(4.85), Inches(4.7), Inches(1.6),
       "Training only on “verified clean” data assumes a labelling capability production does not have. The final model trades ~5 pts of recall for fewer false alarms and zero labelling cost.",
       size=14.5, color=MUTED, italic=True)
    cite(s, "Tables 4.2.1–4.2.2, D_train = 2,171 alerts (1,534 benign / 637 attack) — thesis §4.2")
    add_notes(s, "The training story is the methodological heart. Production logs are never guaranteed clean, "
                 "so I train on everything: IF with contamination auto, AE with three rounds of iterative "
                 "trimming — train, drop the worst-reconstructed quarter, retrain, drop 15%, final fit and "
                 "percentile calibration. Right table: versus the supervised-ish baseline I lose some recall "
                 "but gain precision and cut the FPR to 6.2% — and it needs no labels, which is the honest "
                 "production scenario.")


def s08_results():
    s = new_slide()
    action_title(s, "The ensemble misses the fewest attacks: 41 vs 62 (IF alone) and 208 (AE alone)")
    pic(s, "confusion_light.png", M, Inches(1.5), h=Inches(4.9))
    col_header(s, Inches(7.6), Inches(1.55), Inches(5.1), "On 2,171 real alerts")
    rich_bullets(s, Inches(7.6), Inches(2.05), Inches(5.1), Inches(3.0), [
        ("Ensemble:", "precision 82.1% · recall 93.6% · F1 87.5%"),
        ("Standalone AUC:", "IF 0.942 — the workhorse; AE 0.748 — the second opinion that cuts false alarms."),
        ("Real data, real attacks:", "months of live Wazuh telemetry + a WordPress site attacked in the wild for two months."),
    ], size=16.5, gap=12)
    highlight_box(s, Inches(7.6), Inches(5.3), Inches(5.1), Inches(1.1),
                  "Design intent: fewest missed attacks beats the prettiest F1 — misses are what hurt a SOC", size=14.5)
    cite(s, "Figure 4.2.7 (ensemble confusion matrix), Table 4.2.4, Figure 4.2.6 (ROC) — thesis §4.2")
    add_notes(s, "Results on the full evaluation set. The exhibit is the ensemble confusion matrix: 596 attacks "
                 "caught, 41 missed — versus 62 missed by IF alone and 208 by AE alone. That is the number that "
                 "matters for a security system. Precision 82%, so roughly one in five flags is a false alarm — "
                 "acceptable for triage where the cost asymmetry favours recall. Data is real: my own deployment "
                 "plus a genuinely attacked WordPress site.")


def s09_rag():
    s = new_slide()
    action_title(s, "The RAG copilot grounds a local LLM in a self-updating threat knowledge base — every choice ablated")
    # pipeline: query → embed → dense+sparse → RRF → k=5 → LLM → cited answer
    y0 = Inches(1.85)
    def pbox(x, w, lines, c=NAVY, fill=SKY, tcol=None, h=Inches(1.05)):
        b = box(s, x, y0, w, h, fill, c, 1.3)
        box_text(b, lines, size=13.5, color=tcol or c)
        return b
    pbox(M, Inches(1.55), ["alert /", "analyst query"])
    arrow(s, M + Inches(1.55), y0 + Inches(0.5), Inches(2.5), y0 + Inches(0.5))
    pbox(Inches(2.5), Inches(1.75), ["MiniLM-L6-v2", "384-d embedding"])
    arrow(s, Inches(4.25), y0 + Inches(0.5), Inches(4.7), y0 + Inches(0.5))
    b1 = box(s, Inches(4.7), y0 - Inches(0.42), Inches(2.1), Inches(0.8), SKY, BLUE, 1.2)
    box_text(b1, ["dense · cosine"], size=13.5, color=BLUE)
    b2 = box(s, Inches(4.7), y0 + Inches(0.68), Inches(2.1), Inches(0.8), SKY, BLUE, 1.2)
    box_text(b2, ["sparse · BM25"], size=13.5, color=BLUE)
    arrow(s, Inches(6.8), y0 - Inches(0.02), Inches(7.25), y0 + Inches(0.4))
    arrow(s, Inches(6.8), y0 + Inches(1.08), Inches(7.25), y0 + Inches(0.66))
    pbox(Inches(7.25), Inches(1.5), ["RRF fusion", "top k = 5"])
    arrow(s, Inches(8.75), y0 + Inches(0.5), Inches(9.2), y0 + Inches(0.5))
    lb = box(s, Inches(9.2), y0, Inches(3.45), Inches(1.05), NAVY)
    box_text(lb, ["llama3.2 · grounded prompt", "“answer ONLY from context, cite [Source], refuse if absent”"],
             size=13.5, color=WHITE)
    # design ablations row
    col_header(s, M, Inches(3.5), Inches(12), "Each design decision was measured, not assumed", color=NAVY)
    abl = [
        ("Chunking: 512 tokens + 25% overlap", "beats 128/256/1024 — Recall@5 48.4%, Hit-Rate 82.9%, MRR 0.695"),
        ("k = 5 context chunks", "best hit-rate inside llama3.2's 8K window; k=10+ dilutes attention (lost-in-the-middle)"),
        ("RRF instead of a cross-encoder", "ms-marco re-ranker degraded domain queries; rank fusion needs no learned judgement"),
        ("Self-updating corpus", "scheduled CVE agent: fetch NVD/CISA-KEV → dedup → LLM relevance 0–10 → ≥7 index · 4–6 human queue · <4 drop"),
    ]
    yy = Inches(4.0)
    for head, body in abl:
        tb = s.shapes.add_textbox(M, yy, Inches(12.2), Inches(0.62))
        tf = tb.text_frame; tf.word_wrap = True
        p0 = tf.paragraphs[0]
        r1 = p0.add_run(); r1.text = "▪  " + head + " — "
        r1.font.size = Pt(15); r1.font.bold = True; r1.font.color.rgb = NAVY; r1.font.name = FONT
        r2 = p0.add_run(); r2.text = body
        r2.font.size = Pt(15); r2.font.color.rgb = BODY; r2.font.name = FONT
        yy += Inches(0.62)
    cite(s, "Tables 3.3.2 (chunk ablation), 4.3.2 (top-k), §4.3 stage 5 (RRF vs cross-encoder), CVE agent §4.3 — thesis")
    add_notes(s, "The explanation side. Hybrid retrieval: the same query runs through dense semantic search and "
                 "sparse BM25 keyword search — semantics matches 'SSH brute force' to 'multiple failed "
                 "authentications', BM25 nails exact CVE IDs and rule numbers — fused with reciprocal rank "
                 "fusion. The LLM is forced to answer only from retrieved context with citations, and to refuse "
                 "when the context is empty. Every knob was ablated: chunk size, k, and the re-ranker I removed "
                 "because it hurt. The corpus self-updates through the scheduled CVE agent with a human review "
                 "queue for mid-confidence items.")


def s09b_storage():
    s = new_slide()
    action_title(s, "Two stores, two questions: Qdrant finds what is similar, PostgreSQL records what is true")
    # left = Qdrant, right = PostgreSQL
    cw = Inches(5.85); lx = M; rx = Inches(6.85); hy = Inches(1.55); by = Inches(2.1)
    qh = box(s, lx, hy, cw, Inches(0.52), BLUE)
    box_text(qh, ["Qdrant — vector store · “what is similar?”"], size=16, color=WHITE)
    ph = box(s, rx, hy, cw, Inches(0.52), NAVY)
    box_text(ph, ["PostgreSQL — audit ledger · “what is true?”"], size=16, color=WHITE)
    qbody = box(s, lx, by, cw, Inches(3.55), SKY, BLUE, 1.0, shape=MSO_SHAPE.RECTANGLE)
    pbody = box(s, rx, by, cw, Inches(3.55), SKY, NAVY, 1.0, shape=MSO_SHAPE.RECTANGLE)
    def fill(bx, rows, c):
        tf = bx.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.16); tf.margin_top = Inches(0.14)
        for i, (lead, rest) in enumerate(rows):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_after = Pt(9); p.line_spacing = 1.05
            r = p.add_run(); r.text = "▪ " + lead
            r.font.size = Pt(15); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
            if rest:
                r2 = p.add_run(); r2.text = " " + rest
                r2.font.size = Pt(15); r2.font.color.rgb = BODY; r2.font.name = FONT
    fill(qbody, [
        ("Data:", "text → 384-d embeddings (MiniLM) + BM25 sparse vectors"),
        ("Index:", "HNSW approximate nearest-neighbour, sub-linear"),
        ("Query:", "hybrid dense + sparse, fused with RRF"),
        ("Holds:", "verified threat intel + alert episodes only"),
        ("Answers:", "“which knowledge means the same as this alert?”"),
    ], BLUE)
    fill(pbody, [
        ("Data:", "5 relational tables + raw JSONB snapshots"),
        ("Guarantees:", "ACID, keys, constraints — the cve_decisions ledger"),
        ("Governance:", "medium-confidence CVEs quarantined, never in RAG"),
        ("Durability:", "full original payload survives model/index change"),
        ("Answers:", "“what exactly happened, and is it consistent?”"),
    ], NAVY)
    highlight_box(s, Inches(3.0), Inches(5.9), Inches(7.35), Inches(0.72),
                  "Bridge: cve_decisions.qdrant_id links each audited fact to its searchable vector — a foreign key across paradigms", size=14.5)
    cite(s, "Thesis §3.2 (embedding model, HNSW, PostgreSQL audit helper); schema.sql + scheduled_agent/schema_additions.sql")
    add_notes(s, "A deliberate two-database design, because they answer different questions. Qdrant is the vector "
                 "store: text becomes 384-dimensional embeddings plus sparse BM25 vectors, indexed with HNSW for "
                 "fast approximate nearest-neighbour, queried as a hybrid fused by reciprocal rank fusion — it "
                 "answers 'what is semantically similar'. PostgreSQL is the relational audit layer: five tables "
                 "plus raw JSONB snapshots, with ACID guarantees and constraints — it answers 'what exactly "
                 "happened and is the record consistent'. Similarity search can't run on a B-tree, and audit "
                 "guarantees can't run on a vector index — so I use both. The bridge is one column: cve_decisions "
                 "dot qdrant_id links an audited decision to its vector point — effectively a foreign key that "
                 "crosses the two paradigms. And governance: only verified intelligence reaches Qdrant; "
                 "medium-confidence items are quarantined in Postgres.")


def s09c_rag_retrieval():
    s = new_slide()
    action_title(s, "Retrieval works where it matters: CVEs and APT profiles near-perfect; MITRE techniques are the weak spot")
    pic(s, "rag_retrieval_light.png", M, Inches(1.5), w=Inches(7.55))
    col_header(s, Inches(8.4), Inches(1.55), Inches(4.3), "Over 82 annotated queries")
    rich_bullets(s, Inches(8.4), Inches(2.15), Inches(4.3), Inches(3.3), [
        ("Overall:", "Hit-Rate 82.9% · MRR 0.695 at k=5 — the right source is usually ranked at or near the top."),
        ("Strong:", "CVEs & APT profiles (100%) — distinctive, narrow vocabulary embeds cleanly."),
        ("Weak:", "MITRE techniques (46%) — one query maps to many overlapping sub-techniques."),
    ], size=15.5, gap=11)
    highlight_box(s, Inches(8.4), Inches(5.35), Inches(4.3), Inches(1.15),
                  "Honest finding: a cross-encoder re-ranker HURT domain queries → removed. RRF needs no learned judgement.", size=14)
    cite(s, "Appendix 3–4 (retrieval by category), §4.3 — thesis; hybrid dense+sparse+RRF over all 82 queries")
    add_notes(s, "Before generation, does retrieval even find the right evidence? On 82 human-annotated queries, "
                 "hit-rate 82.9%, MRR 0.695 — the correct source is usually at or near the top. The chart breaks "
                 "it down by category: CVEs and APT profiles are essentially perfect because their vocabulary is "
                 "distinctive. The weak spot is MITRE technique queries at 46% recall — one query legitimately "
                 "maps to many overlapping sub-techniques, so 'the' relevant chunk is ambiguous. The honest "
                 "engineering finding: I tried a cross-encoder re-ranker and it made domain queries worse — a "
                 "model trained on web search doesn't know rule 5710 relates to SSH — so I removed it. RRF just "
                 "trusts rank position from two independent retrievers and needs no learned judgement.")


def s10_rag_results():
    s = new_slide()
    action_title(s, "Generation is verifiably grounded: 98.3% citation precision — and it refuses to invent threats")
    pic(s, "rag_gen_light.png", M, Inches(1.6), w=Inches(7.5))
    col_header(s, Inches(8.5), Inches(1.6), Inches(4.2), "Why it matters")
    rich_bullets(s, Inches(8.5), Inches(2.15), Inches(4.2), Inches(3.4), [
        ("Every claim is auditable:", "inline [Source] tags map to retrieved chunks — 117/119 verified correct."),
        ("Asked about a fake CVE / APT group:", "refused all 6 fabricated-intelligence probes instead of hallucinating."),
        ("Never over-refuses:", "0 of 82 answerable queries wrongly declined — calibrated, not just cautious."),
        ("Retrieval base:", "Hit-Rate 82.9% · MRR 0.695 at k=5 on 82 annotated queries."),
    ], size=15.5, gap=10)
    cite(s, "Figure 4.3.6 (generation evaluation), §4.3 refusal-behaviour set — thesis; RAGAS-style claim decomposition [27]")
    add_notes(s, "Does the copilot actually stay grounded? Citation precision 98.3% — of 119 citations the "
                 "model produced, 117 point at a chunk that really supports the sentence. Faithfulness 0.71 — "
                 "the rest is mostly generic background, not contradictions. The security-critical result: on "
                 "purpose-built unanswerable probes — a fake CVE, an invented APT group — it refused every "
                 "fabricated-intelligence question rather than making threat intel up, and it never refused a "
                 "question it could answer.")


def s10b_cve_agent():
    s = new_slide()
    action_title(s, "The knowledge base updates itself: a bounded agent scores and routes every new CVE — daily, audited")
    AMBER = RGBColor(0xC5, 0x7A, 0x00)
    # ---- top pipeline: Fetch → Dedup → Boost → LLM Score ----
    py = Inches(1.75); ph = Inches(1.18); pw = Inches(2.7); gap = Inches(0.4)
    steps = [
        ("Fetch", ["NVD + CISA-KEV", "since last cursor"], NAVY),
        ("Dedup", ["vs threat_intel", "Postgres ledger"], BLUE),
        ("Boost", ["+ environment evidence", "local rules · alert history"], BLUE),
        ("LLM Score", ["llama3.2 → relevance 0–10", "final = min(base + boost, 10)"], NAVY),
    ]
    xs = []
    x = M
    for i, (head, lines, c) in enumerate(steps):
        b = box(s, x, py, pw, ph, SKY, c, 1.4)
        box_text(b, [head] + lines, size=15, color=c)
        xs.append(x)
        if i > 0:
            arrow(s, xs[i - 1] + pw, py + ph / 2, x, py + ph / 2, MUTED, 2.2)
        x += pw + gap
    llm_cx = xs[3] + pw / 2
    # ---- routing outcomes ----
    ry = Inches(4.35); rh = Inches(1.2); rw = Inches(3.75); rgap = Inches(0.45)
    outs = [
        ("≥ 7   AUTO-INDEX", "indexed → Qdrant + Postgres", GREEN),
        ("4–6   REVIEW QUEUE", "held OUT of RAG · human decides", AMBER),
        ("< 4   DROP + AUDIT", "reasoning logged, never silent", REDX),
    ]
    rx = M
    for head, sub, c in outs:
        b = box(s, rx, ry, rw, rh, WHITE, c, 1.8)
        tf = b.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]; p0.alignment = PP_ALIGN.CENTER
        r = p0.add_run(); r.text = head
        r.font.size = Pt(18); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
        p1 = tf.add_paragraph(); p1.alignment = PP_ALIGN.CENTER
        r2 = p1.add_run(); r2.text = sub
        r2.font.size = Pt(13.5); r2.font.color.rgb = BODY; r2.font.name = FONT
        arrow(s, llm_cx, py + ph, rx + rw / 2, ry, MUTED, 2.0)
        rx += rw + rgap
    # ---- ledger + safety note ----
    highlight_box(s, M, Inches(5.9), CW, Inches(0.6),
                  "Every decision lands in the cve_decisions ledger — reversible, auditable, no silent RAG pollution", size=15)
    tx(s, M, Inches(6.6), CW, Inches(0.4),
       "One bounded LLM step; on failure → conservative score 5 → review queue, never a silent drop.",
       size=14, color=MUTED, italic=True)
    cite(s, "Scheduled CVE ingestion agent — thesis §4.3; Postgres cve_decisions ledger + pending_review view")
    add_notes(s, "This is the third contribution: the knowledge base maintains itself. A scheduled agent runs "
                 "daily with a deliberately fixed control flow — Fetch new CVEs from NVD and CISA-KEV since the "
                 "last cursor, Dedup against the Postgres ledger, Boost relevance with local evidence (my own "
                 "Wazuh rules and alert history), then exactly ONE bounded LLM step that scores relevance 0 to 10. "
                 "The score routes the CVE: 7 or above auto-indexes into Qdrant and Postgres, 4 to 6 is held OUT "
                 "of the RAG and queued for a human, below 4 is dropped but its reasoning is logged. The safety "
                 "properties matter: only one LLM call so it can't spiral, and if the LLM fails it defaults to a "
                 "conservative 5 which routes to human review — never a silent drop. Every decision is written to "
                 "the cve_decisions audit ledger, so the whole thing is reversible and there's no silent pollution "
                 "of the knowledge base.")


def s11_conclusions():
    s = new_slide(dark=True)
    tx(s, M, Inches(0.35), Inches(9), Inches(0.5), "Conclusions", size=20, color=LIGHTB)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(0.88), CW, Emu(36576))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    items = [
        ("1.  Unsupervised detection works on real telemetry:",
         "the IF+AE ensemble reaches recall 93.6% with the fewest missed attacks (41), trained with zero labels on contaminated data."),
        ("2.  Grounded explanation is measurable, not aspirational:",
         "98.3% citation precision, refusal of fabricated threats, 0 over-refusals — the copilot earns analyst trust."),
        ("3.  The knowledge base maintains itself:",
         "a scheduled agent scores and routes new CVEs into the RAG index with a human review queue and full audit ledger in PostgreSQL."),
    ]
    y = Inches(1.15)
    for head, body in items:
        tb = s.shapes.add_textbox(M, y, CW, Inches(1.25))
        tf = tb.text_frame; tf.word_wrap = True
        p0 = tf.paragraphs[0]; p0.line_spacing = 1.1
        r1 = p0.add_run(); r1.text = head + " "
        r1.font.size = Pt(20); r1.font.bold = True; r1.font.color.rgb = WHITE; r1.font.name = FONT
        r2 = p0.add_run(); r2.text = body
        r2.font.size = Pt(19); r2.font.color.rgb = LIGHTB; r2.font.name = FONT
        y += Inches(1.32)
    tx(s, M, Inches(5.25), CW, Inches(0.85),
       "Honest limits: ~2,200 training alerts; per-alert scoring (no UEBA yet); local LLM latency. Next: entity-level behaviour analytics, supervised comparison baseline, domain-tuned re-ranker.",
       size=15, color=LIGHTB, italic=True)
    ln2 = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(6.25), Inches(3.2), Emu(27432))
    ln2.fill.solid(); ln2.fill.fore_color.rgb = BLUE; ln2.line.fill.background()
    tx(s, M, Inches(6.45), CW, Inches(0.7),
       "github.com/Diana-Secarea/advanced-ai-siem  ·  diana.secarea1111@gmail.com  ·  feedback welcome",
       size=15, color=WHITE)
    add_notes(s, "Three takeaways. One: unsupervised anomaly detection is viable on real SIEM telemetry — "
                 "best recall of the three configurations, no labels needed. Two: LLM explanation can be held "
                 "to measurable grounding standards — and passes. Three: the knowledge base stays fresh "
                 "autonomously, with humans only where confidence is medium. Full implementation is public on "
                 "GitHub. I'll stop here — happy to take questions.  [This slide stays on screen during Q&A.]")


def s12_references():
    s = new_slide()
    tx(s, M, Inches(0.22), CW, Inches(0.4), "Appendix — References (selection cited in this deck)",
       size=14, color=MUTED, italic=True)
    tx(s, M, Inches(0.62), CW, Inches(0.55), "References", size=24, color=NAVY, bold=True)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(1.18), CW, Emu(22860))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background()
    refs = [
        "[2] Wazuh Inc., “Wazuh Documentation”, https://documentation.wazuh.com/, 2024.",
        "[4] World Economic Forum, “The Global Risks Report 2024”, Geneva, 2024.",
        "[7] E. & L. S. Okafor, “Unsupervised Learning Framework for Cyber Threat Detection, Anomaly Identification, and Alert Prioritization”, Applied Sciences, 2026.",
        "[20] F. T. Liu, K. M. Ting, Z.-H. Zhou, “Isolation Forest”, IEEE International Conference on Data Mining, 2008.",
        "[23] U. Michelucci, “An Introduction to Autoencoders”, arXiv:2201.03898, 2022.",
        "[25] P. Lewis et al., “Retrieval-Augmented Generation for Knowledge-Intensive NLP Tasks”, NeurIPS 2020.",
        "[27] Y. Gao et al., “Retrieval-Augmented Generation for Large Language Models: A Survey”, arXiv:2312.10997, 2023.",
        "[29] N. Reimers, I. Gurevych, “Sentence-BERT: Sentence Embeddings using Siamese BERT-Networks”, EMNLP 2019.",
        "[31] Qdrant Solutions GmbH, “Qdrant Documentation”, 2026.",
        "[32] D. M. Secărea, “AI-SIEM Threat Engine”, https://github.com/Diana-Secarea/advanced-ai-siem, 2026.",
    ]
    tb = s.shapes.add_textbox(M, Inches(1.4), CW, Inches(5.5))
    tf = tb.text_frame; tf.word_wrap = True
    for i, rtext in enumerate(refs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run(); r.text = rtext
        r.font.size = Pt(13.5); r.font.color.rgb = BODY; r.font.name = FONT
    tx(s, M, Inches(6.95), CW, Inches(0.4),
       "Full bibliography (32 entries): thesis pp. 44–46.", size=12, color=MUTED)
    add_notes(s, "Appendix slide — navigate here only if a reference question comes up.")


# ============================ BUILD ============================
if __name__ == "__main__":
    chart_ablation()
    chart_confusion()
    chart_rag_generation()
    chart_rag_retrieval()
    s01_title()
    s02_motivation()
    s03_system()
    s04_technologies()
    s05_ensemble()
    s06_features()
    s07_training()
    s08_results()
    s09_rag()
    s09b_storage()
    s09c_rag_retrieval()
    s10_rag_results()
    s10b_cve_agent()
    s11_conclusions()
    s12_references()
    out = os.path.join(HERE, "GATA_Defense_Presentation.pptx")
    prs.save(out)
    print("saved:", out, f"({len(prs.slides._sldIdLst)} slides)")
