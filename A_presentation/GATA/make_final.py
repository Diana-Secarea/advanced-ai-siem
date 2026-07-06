#!/usr/bin/env python3
"""Build the FINAL short defense deck — 9 slides, committee format.

Structure (per committee rules: minimalist, no animations, high contrast,
max 14-16 words per bullet, one Demo/Q&A slide):
  1 Title · 2 Objectives (what) · 3 Contents · 4 Motivation (why)
  5 Architecture + tech stack · 6 ML processing flow · 7 Results
  8 Conclusions & future work · 9 Demo + Q&A

Same academic-pptx style as make_gata.py (white slides, action titles,
navy/blue palette, one exhibit per slide, in-slide citations).

Run:  ../../ai_threat_engine_starter/venv/bin/python make_final.py
Outputs: final_presentation.pptx  (reuses charts already in img/)
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
IMG = os.path.join(HERE, "img")

# ---------- academic palette (identical to make_gata.py) ----------
NAVY   = RGBColor(0x1F, 0x4E, 0x79)
BLUE   = RGBColor(0x2E, 0x75, 0xB6)
SKY    = RGBColor(0xEB, 0xF3, 0xFA)
BODY   = RGBColor(0x2D, 0x2D, 0x2D)
MUTED  = RGBColor(0x77, 0x77, 0x77)
RULE   = RGBColor(0xCC, 0xCC, 0xCC)
HL_BG  = RGBColor(0xFF, 0xF2, 0xCC)
HL_BR  = RGBColor(0xE6, 0xC8, 0x00)
HL_TX  = RGBColor(0x7A, 0x52, 0x00)
GREEN  = RGBColor(0x2E, 0x7D, 0x5B)
REDX   = RGBColor(0xB0, 0x35, 0x2F)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
LIGHTB = RGBColor(0xA0, 0xBB, 0xDD)
FONT = "Calibri"

SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
M = Inches(0.6)
CW = Inches(12.13)

prs = Presentation()
prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
BLANK = prs.slide_layouts[6]


# ---------- helpers (identical style to make_gata.py) ----------

def new_slide(dark=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = NAVY if dark else WHITE
    return s


def add_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def tx(slide, x, y, w, h, text, size=20, color=BODY, bold=False, italic=False,
       align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    lines = text.split("\n") if isinstance(text, str) else text
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = spacing
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
        r.font.color.rgb = color; r.font.name = FONT
    return tb


def action_title(slide, text, size=25):
    tx(slide, M, Inches(0.28), CW, Inches(1.0), text, size=size, color=NAVY, bold=True)
    ln = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, M, Inches(1.26), CW, Emu(22860))
    ln.fill.solid(); ln.fill.fore_color.rgb = RULE; ln.line.fill.background()


def rich_bullets(slide, x, y, w, h, items, size=19, gap=10, marker_color=BLUE):
    """items: list of (bold_lead, rest) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = True
    for i, (lead, rest) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.line_spacing = 1.08
        rm = p.add_run(); rm.text = "▪  "
        rm.font.size = Pt(size); rm.font.color.rgb = marker_color; rm.font.name = FONT
        r1 = p.add_run(); r1.text = lead
        r1.font.size = Pt(size); r1.font.bold = True; r1.font.color.rgb = BODY; r1.font.name = FONT
        if rest:
            r2 = p.add_run(); r2.text = " " + rest
            r2.font.size = Pt(size); r2.font.color.rgb = BODY; r2.font.name = FONT
    return tb


def cite(slide, text, y=Inches(7.02)):
    tx(slide, M, y, CW, Inches(0.36), text, size=12, color=MUTED)


def box(slide, x, y, w, h, fill, line_color=None, line_w=1.0,
        shape=MSO_SHAPE.ROUNDED_RECTANGLE):
    sh = slide.shapes.add_shape(shape, x, y, w, h)
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if line_color is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line_color; sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    try:
        sh.adjustments[0] = 0.12
    except Exception:
        pass
    return sh


def box_text(sh, lines, size=14, color=BODY, bold_first=True):
    tf = sh.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.03)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run(); r.text = line
        r.font.size = Pt(size if i == 0 else size - 2)
        r.font.bold = bold_first and i == 0
        r.font.color.rgb = color; r.font.name = FONT


def arrow(slide, x1, y1, x2, y2, color=MUTED, w=2.0):
    conn = slide.shapes.add_connector(2, x1, y1, x2, y2)
    conn.line.color.rgb = color; conn.line.width = Pt(w)
    le = conn.line._get_or_add_ln()
    tail = le.makeelement(qn("a:tailEnd"), {"type": "triangle", "w": "med", "len": "med"})
    le.append(tail)


def highlight_box(slide, x, y, w, h, text, size=14):
    sh = box(slide, x, y, w, h, HL_BG, HL_BR, 1.0)
    box_text(sh, [text], size=size, color=HL_TX, bold_first=True)
    return sh


def pic(slide, path, x, y, w=None, h=None):
    return slide.shapes.add_picture(os.path.join(IMG, path), x, y, w, h)


# ============================ SLIDES ============================

def s1_title():
    s = new_slide(dark=True)
    tx(s, Inches(0.9), Inches(0.75), Inches(11.5), Inches(0.5),
       "BACHELOR THESIS DEFENSE", size=15, color=LIGHTB)
    tx(s, Inches(0.9), Inches(1.5), Inches(11.6), Inches(2.6),
       "Architecture and Evaluation of a SIEM\nIntegrating Unsupervised Anomaly Detection\nand Knowledge-Grounded Language Models",
       size=33, color=WHITE, bold=True, spacing=1.08)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.9), Inches(4.35), Inches(2.2), Emu(45720))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    tx(s, Inches(0.9), Inches(4.6), Inches(11), Inches(0.45),
       "Secărea Diana Maria", size=21, color=WHITE, bold=True)
    tx(s, Inches(0.9), Inches(5.12), Inches(11.6), Inches(1.0),
       "Coordinator: prof. univ. dr. Cătălin Boja\nBucharest University of Economic Studies · Faculty of Cybernetics, Statistics and Economic Informatics",
       size=15, color=LIGHTB)
    tx(s, Inches(0.9), Inches(6.45), Inches(11), Inches(0.4),
       "Economic Informatics Program · Bucharest 2026", size=13, color=LIGHTB)
    add_notes(s, "Good morning. My thesis integrates two AI paradigms into the Wazuh SIEM: "
                 "unsupervised anomaly detection and a knowledge-grounded language model. "
                 "Everything is implemented, deployed and evaluated — and I will show it live.")


def s2_objectives():
    s = new_slide()
    action_title(s, "Objective — score every alert, explain every score, keep knowledge fresh")
    # centered objective paragraph
    tx(s, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.0),
       "Extend the open-source Wazuh SIEM with two cooperating AI paradigms,\nimplemented, deployed, and evaluated on real security telemetry.",
       size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER, spacing=1.15)
    rich_bullets(s, Inches(2.0), Inches(3.3), Inches(9.6), Inches(2.8), [
        ("Detect:", "score every alert 0–100 with an unsupervised Isolation Forest + Autoencoder ensemble."),
        ("Explain:", "a local LLM answers only from cited, retrieved threat intelligence."),
        ("Stay current:", "a scheduled agent ingests new CVEs autonomously, with full audit."),
    ], size=20, gap=16)
    highlight_box(s, Inches(3.4), Inches(6.15), Inches(6.5), Inches(0.6),
                  "Fully local — security data never leaves the machine", size=15)
    cite(s, "Thesis §1 — objectives and scope")
    add_notes(s, "One sentence: I extend Wazuh with two AI paradigms. Three concrete objectives: "
                 "detect with an unsupervised ensemble, explain with a grounded local LLM, "
                 "and keep the knowledge base current autonomously. Everything runs locally.")


def s3_contents():
    s = new_slide()
    action_title(s, "Contents — what was analyzed, discovered, proposed, and evaluated")
    items = [
        ("Analyzed", "where rule-based SIEMs fail: alert fatigue, unknown attacks, missing context", "thesis §2", BLUE),
        ("Discovered", "two unsupervised detectors fail differently — fusing them cuts missed attacks", "thesis §3.2", NAVY),
        ("Proposed", "an AI Threat-Hunting Engine: ensemble scoring, RAG copilot, CVE agent", "thesis §3", GREEN),
        ("Evaluated", "2,171 real alerts · 82 annotated queries · one live deployment", "thesis §4", BLUE),
    ]
    y = Inches(1.75)
    for head, body, ref, c in items:
        b = box(s, Inches(1.1), y, Inches(11.1), Inches(1.05), SKY, c, 1.4)
        tf = b.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.18); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        r = p1.add_run(); r.text = head + "   "
        r.font.size = Pt(20); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
        r2 = p1.add_run(); r2.text = body
        r2.font.size = Pt(17); r2.font.color.rgb = BODY; r2.font.name = FONT
        r3 = p1.add_run(); r3.text = "   — " + ref
        r3.font.size = Pt(13); r3.font.color.rgb = MUTED; r3.font.name = FONT
        y += Inches(1.25)
    add_notes(s, "The talk follows the thesis structure: I analyzed the limits of rule-based SIEMs, "
                 "discovered that two unsupervised models fail in complementary ways, proposed the "
                 "AI Threat-Hunting Engine architecture, and evaluated it on real data end to end.")


def s4_motivation():
    s = new_slide()
    action_title(s, "Why — analysts drown in alerts while novel attacks pass undetected")
    rich_bullets(s, M, Inches(1.65), Inches(7.1), Inches(3.6), [
        ("Alert fatigue:", "SOCs average ~4,490 alerts per day; about 67% are never investigated."),
        ("Blind to the unknown:", "signature rules match only known patterns — zero-days pass silently."),
        ("No context:", "each alert arrives bare; every investigation starts from zero."),
    ], size=20, gap=16)
    qb = box(s, Inches(8.1), Inches(1.75), Inches(4.55), Inches(3.1), SKY, BLUE, 1.5)
    box_text(qb, ["Research question",
                  "",
                  "Can unsupervised scoring",
                  "+ grounded reasoning",
                  "cut triage time and earn",
                  "analyst trust — on a real,",
                  "running deployment?"], size=17, color=NAVY)
    highlight_box(s, Inches(1.6), Inches(5.7), Inches(10.1), Inches(0.62),
                  "How it helps: fewer missed attacks · faster triage · explanations analysts can verify", size=15)
    cite(s, "Industry SOC statistics [7]; WEF Global Risks Report 2024 [4] — thesis §2")
    add_notes(s, "Why this matters: most alerts are never investigated, signature detection cannot see "
                 "novel attacks, and alerts carry no context. The research question: can unsupervised "
                 "scoring plus grounded reasoning fix triage on a real deployment, not a benchmark.")


def s5_architecture():
    s = new_slide()
    action_title(s, "How — architecture: three AI layers on top of standard Wazuh")
    p = pic(s, "architecture_thesis.png", M, Inches(1.45), w=Inches(7.3))
    x = Inches(8.2); w = Inches(4.5)
    labels = [
        ("1 · Anomaly ensemble", "IF + Autoencoder score every alert, fused into confidence tiers", NAVY),
        ("2 · RAG copilot", "local LLM grounded in MITRE, YARA, Sigma, CVEs — cited answers", BLUE),
        ("3 · CVE agent", "fetches, scores, routes new CVEs; PostgreSQL audit ledger", GREEN),
    ]
    y = Inches(1.6)
    for head, body, c in labels:
        b = box(s, x, y, w, Inches(1.05), SKY, c, 1.4)
        tf = b.text_frame; tf.word_wrap = True
        tf.margin_left = Inches(0.12); tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p1 = tf.paragraphs[0]
        r = p1.add_run(); r.text = head
        r.font.size = Pt(16); r.font.bold = True; r.font.color.rgb = c; r.font.name = FONT
        p2 = tf.add_paragraph()
        r2 = p2.add_run(); r2.text = body
        r2.font.size = Pt(12.5); r2.font.color.rgb = BODY; r2.font.name = FONT
        y += Inches(1.2)
    # technology stack chips (3 rows, kept clear of the citation line)
    tx(s, x, Inches(5.18), w, Inches(0.3), "TECHNOLOGY STACK", size=13, color=MUTED, bold=True)
    stack = ["Wazuh 4.14", "Python 3.13", "scikit-learn", "Ollama · llama3.2",
             "Qdrant", "PostgreSQL 17", "Flask", "Docker"]
    cx, cy = x, Inches(5.52)
    chip_h = Inches(0.4)
    widths = [Inches(1.35), Inches(1.35), Inches(1.4), Inches(1.85),
              Inches(1.0), Inches(1.6), Inches(0.85), Inches(1.0)]
    row_x = cx
    for name, wd in zip(stack, widths):
        if row_x + wd > x + w:
            row_x = cx; cy = cy + Inches(0.48)
        ch = box(s, row_x, cy, wd, chip_h, WHITE, BLUE, 1.0)
        box_text(ch, [name], size=12, color=NAVY, bold_first=True)
        row_x = row_x + wd + Inches(0.12)
    cite(s, "Figure 3.1.1 — thesis, p. 9. Endpoint & manager layers = Wazuh; engine, storage, visualisation = this work.")
    add_notes(s, "The left half is standard Wazuh — agents, manager, rule engine. Everything above is my "
                 "contribution: the ensemble scores, the RAG copilot explains with citations, the CVE agent "
                 "keeps knowledge fresh. Stack is fully local and open source: scikit-learn, Ollama with "
                 "llama3.2, Qdrant, PostgreSQL, Docker — no cloud API, no per-token cost.")


def s6_ml_flow():
    s = new_slide()
    action_title(s, "How — the ML flow: train on raw alerts, score every new alert in ~12 ms")
    # ---- TRAINING lane ----
    tx(s, M, Inches(1.55), Inches(5.5), Inches(0.35), "TRAINING — offline, zero labels",
       size=14, color=NAVY, bold=True)
    train_boxes = [
        ("2,171 real alerts", "months of live telemetry", Inches(2.5)),
        ("16 features", "per alert, z-scaled", Inches(2.1)),
        ("Iterative trimming ×3", "drop worst-reconstructed 25% → 15%", Inches(3.1)),
        ("IF + AE models", "calibrated 0–100, saved", Inches(2.5)),
    ]
    bx = M; by = Inches(2.0); bh = Inches(0.95)
    lane1 = []
    for head, sub, bw in train_boxes:
        b = box(s, bx, by, bw, bh, SKY, NAVY, 1.2)
        box_text(b, [head, sub], size=14, color=NAVY)
        lane1.append((bx, bw))
        bx = bx + bw + Inches(0.55)
    for i in range(len(lane1) - 1):
        x_end = lane1[i][0] + lane1[i][1]
        arrow(s, x_end + Inches(0.06), by + bh / 2, lane1[i + 1][0] - Inches(0.06), by + bh / 2, NAVY, 2.2)
    # ---- SCORING lane ----
    tx(s, M, Inches(3.85), Inches(4.0), Inches(0.35), "SCORING — live, per alert",
       size=14, color=GREEN, bold=True)
    score_boxes = [
        ("New Wazuh alert", "from alerts.json", Inches(2.2)),
        ("Same 16 features", "identical extractor", Inches(2.2)),
        ("IF score · AE score", "two independent opinions", Inches(2.4)),
        ("Fusion", "0.45·IF + 0.55·AE", Inches(1.7)),
        ("Tier", "CRITICAL · HIGH · POSSIBLE · NORMAL", Inches(2.5)),
    ]
    bx = M; by2 = Inches(4.3)
    lane2 = []
    for head, sub, bw in score_boxes:
        b = box(s, bx, by2, bw, bh, WHITE, GREEN, 1.4)
        box_text(b, [head, sub], size=13, color=BODY)
        lane2.append((bx, bw))
        bx = bx + bw + Inches(0.28)
    for i in range(len(lane2) - 1):
        x_end = lane2[i][0] + lane2[i][1]
        arrow(s, x_end + Inches(0.04), by2 + bh / 2, lane2[i + 1][0] - Inches(0.04), by2 + bh / 2, GREEN, 2.2)
    # connector: trained models feed the live IF·AE scoring box
    models_bottom_x = lane1[3][0] + train_boxes[3][2] / 2
    ifae_top_x = lane2[2][0] + score_boxes[2][2] / 2
    arrow(s, models_bottom_x, by + bh + Inches(0.05), ifae_top_x, by2 - Inches(0.05), NAVY, 2.2)
    highlight_box(s, Inches(1.6), Inches(5.85), Inches(10.1), Inches(0.62),
                  "No labels anywhere: the model bootstraps its own clean set from contaminated data", size=15)
    cite(s, "Training: thesis §3.2, Tables 4.2.1–4.2.2 · Fusion and tiers: thesis §3.2")
    add_notes(s, "What the architecture picture cannot show: how training happens. Top lane, offline: "
                 "months of raw alerts become 16 features; iterative trimming drops the worst-reconstructed "
                 "tail — where attacks hide — and retrains, three rounds, zero labels. Bottom lane, live: "
                 "every new alert goes through the same features, both models score it, scores fuse into a tier. "
                 "About 12 milliseconds per alert, before any human looks at it.")


def s7_results():
    s = new_slide()
    action_title(s, "Results — fewest missed attacks, and a copilot that cites instead of inventing")
    pic(s, "confusion_light.png", M, Inches(1.6), w=Inches(5.6))
    tx(s, M, Inches(6.15), Inches(5.6), Inches(0.4),
       "Ensemble confusion matrix — 2,171 real alerts", size=13, color=MUTED, align=PP_ALIGN.CENTER)
    rich_bullets(s, Inches(6.7), Inches(1.7), Inches(6.0), Inches(4.3), [
        ("Detection:", "recall 93.6% with zero labels; false-positive rate down to 6.2%."),
        ("Fewest misses:", "41 attacks missed — versus 62 (IF alone) and 208 (AE alone)."),
        ("Retrieval:", "Hit-Rate 82.9%, MRR 0.695 across 82 annotated queries."),
        ("Grounding:", "98.3% citation precision — 117 of 119 citations verified correct."),
        ("No hallucinations:", "refused all six fabricated-threat probes; zero wrongful refusals."),
    ], size=17, gap=10)
    highlight_box(s, Inches(6.7), Inches(6.25), Inches(6.0), Inches(0.55),
                  "The copilot cites its evidence — and refuses to invent threats", size=13)
    cite(s, "Figure 4.2.7, Table 4.2.4 (detection) · Figure 4.3.6, §4.3 refusal set (generation) — thesis §4")
    add_notes(s, "Two result families. Detection: the ensemble misses only 41 attacks — the number that "
                 "matters for a SOC — versus 62 and 208 for the single models, at 6.2% false positives, "
                 "trained with zero labels. Generation: 98.3% citation precision, and when probed with a "
                 "fake CVE and an invented APT group it refused every time — while never refusing a real question.")


def s8_conclusions():
    s = new_slide()
    action_title(s, "Conclusions — the three claims hold; limits are stated, each one is future work")
    tx(s, M, Inches(1.6), Inches(6.0), Inches(0.4), "CONCLUSIONS", size=15, color=NAVY, bold=True)
    rich_bullets(s, M, Inches(2.05), Inches(6.1), Inches(3.6), [
        ("Detection works:", "unsupervised ensemble on real telemetry — no labelling required."),
        ("Explanation is measurable:", "grounded, cited, and it passes the tests."),
        ("Knowledge stays current:", "autonomous agent, human review only where uncertain."),
    ], size=18, gap=14)
    tx(s, Inches(7.1), Inches(1.6), Inches(5.6), Inches(0.4), "FUTURE WORK", size=15, color=GREEN, bold=True)
    rich_bullets(s, Inches(7.1), Inches(2.05), Inches(5.6), Inches(3.6), [
        ("UEBA:", "entity-level behaviour analytics across alerts."),
        ("Baseline:", "supervised comparison on hand-labelled data."),
        ("Retrieval:", "domain-tuned re-ranker for MITRE technique queries."),
    ], size=18, gap=14, marker_color=GREEN)
    highlight_box(s, Inches(2.4), Inches(5.75), Inches(8.5), Inches(0.62),
                  "Honest limits: ~2,200 training alerts · per-alert scoring · local-LLM latency", size=14)
    tx(s, M, Inches(6.55), CW, Inches(0.4),
       "github.com/Diana-Secarea/advanced-ai-siem", size=14, color=BLUE, bold=True, align=PP_ALIGN.CENTER)
    add_notes(s, "Three conclusions, each shown with evidence: detection works without labels, grounded "
                 "explanation is measurable and passes, and the knowledge base maintains itself. I state "
                 "the limits honestly — corpus size, per-alert scoring, LLM latency — and each maps to a "
                 "concrete future-work item. The full system is public on GitHub.")


def s9_demo():
    s = new_slide(dark=True)
    tx(s, Inches(0.9), Inches(1.9), Inches(11.5), Inches(1.2),
       "Live Demo", size=48, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tx(s, Inches(0.9), Inches(3.15), Inches(11.5), Inches(0.6),
       "real attack  →  live detection  →  grounded, cited explanation",
       size=21, color=LIGHTB, align=PP_ALIGN.CENTER)
    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.57), Inches(4.05), Inches(2.2), Emu(45720))
    ln.fill.solid(); ln.fill.fore_color.rgb = BLUE; ln.line.fill.background()
    tx(s, Inches(0.9), Inches(4.45), Inches(11.5), Inches(0.9),
       "Questions?", size=30, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tx(s, Inches(0.9), Inches(5.75), Inches(11.5), Inches(0.8),
       "Secărea Diana Maria  ·  diana.secarea1111@gmail.com\ngithub.com/Diana-Secarea/advanced-ai-siem",
       size=15, color=LIGHTB, align=PP_ALIGN.CENTER, spacing=1.3)
    add_notes(s, "Switch to the browser: Populate Map, Test the Shield — 14 real attack families through "
                 "the real Wazuh pipeline — click the CRITICAL alert, Investigate with AI, click a source "
                 "chip to show the exact cited passage. Then take questions with this slide up.")


# ============================ BUILD ============================
s1_title()
s2_objectives()
s3_contents()
s4_motivation()
s5_architecture()
s6_ml_flow()
s7_results()
s8_conclusions()
s9_demo()

out = os.path.join(HERE, "final_presentation.pptx")
prs.save(out)
print(f"Saved {out} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")
